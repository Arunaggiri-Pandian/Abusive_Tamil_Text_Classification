"""
Create PowerPoint presentation for Abusive Tamil Text Detection.
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pathlib import Path

OUTPUT_DIR = Path(__file__).parent.parent / "outputs"
FIGURES_DIR = OUTPUT_DIR / "figures"


def add_title_slide(prs, title, subtitle):
    """Add a title slide."""
    slide_layout = prs.slide_layouts[6]  # Blank
    slide = prs.slides.add_slide(slide_layout)

    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(9), Inches(1))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(44)
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER

    # Subtitle
    sub_box = slide.shapes.add_textbox(Inches(0.5), Inches(3.7), Inches(9), Inches(1))
    tf = sub_box.text_frame
    p = tf.paragraphs[0]
    p.text = subtitle
    p.font.size = Pt(24)
    p.alignment = PP_ALIGN.CENTER
    p.font.color.rgb = RGBColor(100, 100, 100)


def add_content_slide(prs, title, content_lines):
    """Add a content slide with bullet points."""
    slide_layout = prs.slide_layouts[6]  # Blank
    slide = prs.slides.add_slide(slide_layout)

    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(32)
    p.font.bold = True

    # Content
    content_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.3), Inches(9), Inches(5.5))
    tf = content_box.text_frame

    for i, line in enumerate(content_lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = line
        p.font.size = Pt(20)
        p.space_after = Pt(10)


def add_table_slide(prs, title, headers, rows):
    """Add a slide with a table."""
    slide_layout = prs.slide_layouts[6]  # Blank
    slide = prs.slides.add_slide(slide_layout)

    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(32)
    p.font.bold = True

    # Table
    cols = len(headers)
    num_rows = len(rows) + 1
    table = slide.shapes.add_table(num_rows, cols, Inches(0.5), Inches(1.3), Inches(9), Inches(0.5 * num_rows)).table

    # Header
    for j, header in enumerate(headers):
        cell = table.cell(0, j)
        cell.text = header
        cell.text_frame.paragraphs[0].font.bold = True
        cell.text_frame.paragraphs[0].font.size = Pt(16)
        cell.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

    # Data rows
    for i, row in enumerate(rows):
        for j, val in enumerate(row):
            cell = table.cell(i + 1, j)
            cell.text = str(val)
            cell.text_frame.paragraphs[0].font.size = Pt(14)
            cell.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER


def add_image_slide(prs, title, image_path, caption=None):
    """Add a slide with an image."""
    slide_layout = prs.slide_layouts[6]  # Blank
    slide = prs.slides.add_slide(slide_layout)

    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(9), Inches(0.6))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(28)
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER

    # Image
    if Path(image_path).exists():
        slide.shapes.add_picture(str(image_path), Inches(1.0), Inches(0.9), width=Inches(8))

    # Caption
    if caption:
        cap_box = slide.shapes.add_textbox(Inches(0.5), Inches(6.5), Inches(9), Inches(0.5))
        tf = cap_box.text_frame
        p = tf.paragraphs[0]
        p.text = caption
        p.font.size = Pt(14)
        p.font.italic = True
        p.alignment = PP_ALIGN.CENTER
        p.font.color.rgb = RGBColor(100, 100, 100)


def create_presentation():
    """Create the full presentation."""
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    # Slide 1: Title
    add_title_slide(prs,
        "Abusive Tamil Text Detection",
        "DravidianLangTech @ ACL 2026 | Team CHMOD_777")

    # Slide 2: Task Overview
    add_content_slide(prs, "Task Overview", [
        "Goal: Detect abusive Tamil text targeting women on social media",
        "",
        "Task Type: Binary classification",
        "  • Non-Abusive vs Abusive",
        "",
        "Dataset: 3,286 train / 366 dev / 913 test samples",
        "",
        "Class Balance: Nearly balanced (51.6% vs 48.4%)",
        "",
        "Evaluation Metric: Macro F1 Score"
    ])

    # Slide 3: Dataset Distribution
    add_image_slide(prs, "Dataset Distribution",
        FIGURES_DIR / "dataset_distribution.png",
        "Nearly balanced classes - No augmentation required")

    # Slide 4: Model Comparison
    add_table_slide(prs, "Model Comparison",
        ["Model", "Max Length", "Best Epoch", "Macro F1"],
        [
            ["MuRIL v2", "256", "9", "82.76% ✓"],
            ["MuRIL v1", "128", "22", "82.50%"],
            ["MuRIL (tuned)", "128", "22", "82.45%"],
            ["XLM-RoBERTa", "256", "10", "81.95%"],
            ["IndicBERT-v3", "256", "39", "74.02%"],
        ])

    # Slide 5: Model Comparison Chart
    add_image_slide(prs, "Model Comparison",
        FIGURES_DIR / "model_comparison.png",
        "MuRIL v2 (256 tokens) achieves best 82.76% Macro F1")

    # Slide 6: Context Length
    add_image_slide(prs, "Context Length Impact",
        FIGURES_DIR / "context_length_comparison.png",
        "Longer context improves both F1 and convergence speed")

    # Slide 7: Per-Class Performance
    add_image_slide(prs, "Per-Class Performance",
        FIGURES_DIR / "per_class_f1.png",
        "Balanced performance across both classes")

    # Slide 8: Confusion Matrix
    add_image_slide(prs, "Confusion Matrix",
        FIGURES_DIR / "confusion_matrix.png",
        "Low confusion between classes")

    # Slide 9: Key Insights
    add_content_slide(prs, "Key Insights", [
        "1. MuRIL outperforms multilingual models:",
        "   • 82.76% vs 81.95% (XLM-RoBERTa)",
        "",
        "2. Language-specific pre-training is crucial",
        "   • MuRIL trained on 17 Indian languages",
        "",
        "3. Longer context helps: 256 > 128 tokens",
        "   • +0.26% F1 improvement",
        "",
        "4. Faster convergence with longer context",
        "   • Epoch 9 vs 22 for best checkpoint"
    ])

    # Slide 10: Submissions
    add_content_slide(prs, "Final Submissions", [
        "Run 1: MuRIL v2 (256 tokens)",
        "  • Macro F1: 82.76% (Best)",
        "",
        "Run 2: MuRIL v1 (128 tokens)",
        "  • Macro F1: 82.50%",
        "",
        "Run 3: XLM-RoBERTa",
        "  • Macro F1: 81.95% (architecture diversity)"
    ])

    # Slide 11: Thank You
    add_title_slide(prs,
        "Thank You!",
        "Team CHMOD_777 | DravidianLangTech @ ACL 2026")

    # Save
    output_path = OUTPUT_DIR / "CHMOD_777_Presentation.pptx"
    prs.save(output_path)
    print(f"Presentation saved to: {output_path}")


if __name__ == "__main__":
    create_presentation()
